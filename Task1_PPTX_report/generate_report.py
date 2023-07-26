import json
import argparse
import collections.abc

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import numpy as np
from io import BytesIO


def create_title_slide(prs, title, content):
    slide_layout = prs.slide_layouts[0]  # Title slide layout index is 0
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]  # Subtitle placeholder index is 1
    title_shape.text = title
    subtitle_shape.text = content

def create_text_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Text slide layout index is 1
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]  # Body text placeholder index is 1
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = content




def create_list_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content with a list layout index is 1
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Create a text frame for the content (the list)
    tf = slide.placeholders[1].text_frame

    for item in content:
        level = item.get('level', 0)
        text = item.get('text', '')
        p = tf.add_paragraph()
        p.text = text
        p.level = level

        

def create_picture_slide(prs, title, image_file):
    slide_layout = prs.slide_layouts[5]  # Title and Content layout index is 5
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Add the picture to the content placeholder
    left = Inches(2)
    top = Inches(2)
    slide.shapes.add_picture(image_file, left, top)



def create_plot_slide(prs, title, data_file, x_label, y_label):
    slide_layout = prs.slide_layouts[0]  # Title and Content layout index is 5
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title

    # Load data from the .dat file using NumPy
    data = np.loadtxt(data_file, unpack=True)

    # Create a figure and plot the data
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots()
    ax.plot(data[0], data[1])
    ax.set_xlabel(x_label)
    ax.set_ylabel(y_label)

    # Save the plot to a BytesIO object (in-memory binary stream)
    buffer = BytesIO()
    plt.savefig(buffer, format="png", bbox_inches="tight")
    buffer.seek(0)

    # Add the plot image to the content placeholder
    left = Inches(1.5)
    top = Inches(2)
    width = Inches(6)
    pic = slide.shapes.add_picture(buffer, left, top, width=width)
    buffer.close()

    # Adjust the plot image size
    pic.width = width
    pic.height = Inches(4.5)  # You can adjust the height as desired

    # Add x-label and y-label to the content placeholder
    txBox = slide.placeholders[1].text_frame
    p = txBox.add_paragraph()
    p.text = x_label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)

    p = txBox.add_paragraph()
    p.text = y_label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)

    # Set font size for the axis labels
    ax.xaxis.label.set_size(12)
    ax.yaxis.label.set_size(12)

    plt.close()




def create_presentation_from_json(json_path):
    with open(json_path) as json_file:
        json_data = json_file.read()

    prs = Presentation()
    presentation_data = json.loads(json_data)

    for slide_data in presentation_data["presentation"]:
        slide_type = slide_data["type"]
        title = slide_data["title"]
        content = slide_data["content"]

        if slide_type == "title":
            create_title_slide(prs, title, content)
        elif slide_type == "text":
            create_text_slide(prs, title, content)
        elif slide_type == "list":
            create_list_slide(prs, title, content)
        elif slide_type == "picture":
            create_picture_slide(prs, title, content)
        elif slide_type == "plot":
            configuration = slide_data.get("configuration", {})
            x_label = configuration.get("x-label", "")
            y_label = configuration.get("y-label", "")
            create_plot_slide(prs, title, content, x_label, y_label)

    return prs

def main():
    parser = argparse.ArgumentParser(description="Create a PowerPoint presentation from a JSON file.")
    parser.add_argument("json_path", help="Path to the JSON file containing presentation data.")
    args = parser.parse_args()

    presentation = create_presentation_from_json(args.json_path)
    output_file = args.json_path.replace(".json", ".pptx")
    presentation.save(output_file)
    print(f"Presentation created and saved as '{output_file}'.")

if __name__ == "__main__":
    main()
